import os
import io
import json
import hashlib
from pptx import Presentation
import pdfplumber
import fitz  # PyMuPDF
import pytesseract
from PIL import Image
from concurrent.futures import ThreadPoolExecutor, as_completed
from .config import TESSERACT_CMD, CACHE_DIR, ENABLE_CACHE, MIN_IMAGE_SIZE, OCR_CONFIG

# Configure Tesseract
pytesseract.pytesseract.tesseract_cmd = TESSERACT_CMD


def _get_file_hash(file_path):
    """Generate hash for file content caching."""
    hasher = hashlib.md5()
    try:
        with open(file_path, 'rb') as f:
            # Read first 64KB for quick hash
            hasher.update(f.read(65536))
        # Include file size and modification time
        stat = os.stat(file_path)
        hasher.update(f"{stat.st_size}_{stat.st_mtime}".encode())
    except:
        hasher.update(file_path.encode())
    return hasher.hexdigest()

def _get_cached_extraction(file_path):
    """Get cached extraction result if available."""
    if not ENABLE_CACHE:
        return None
    cache_key = _get_file_hash(file_path)
    cache_path = os.path.join(CACHE_DIR, f"extract_{cache_key}.json")
    if os.path.exists(cache_path):
        try:
            with open(cache_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
                print(f"  [CACHE HIT] Using cached extraction for {os.path.basename(file_path)}")
                return data.get('content', '')
        except:
            pass
    return None

def _save_extraction_cache(file_path, content):
    """Save extraction result to cache."""
    if not ENABLE_CACHE:
        return
    cache_key = _get_file_hash(file_path)
    cache_path = os.path.join(CACHE_DIR, f"extract_{cache_key}.json")
    try:
        with open(cache_path, 'w', encoding='utf-8') as f:
            json.dump({'content': content, 'source': file_path}, f)
    except:
        pass


# ============= OPTIMIZED OCR =============

def _should_ocr_image(image):
    """Check if image is worth OCR processing."""
    # Skip tiny images (likely icons, logos, bullets)
    if image.width < MIN_IMAGE_SIZE or image.height < MIN_IMAGE_SIZE:
        return False
    # Skip very narrow images (likely decorative lines)
    if image.width < 50 or image.height < 50:
        return False
    # Skip if aspect ratio is extreme (likely decorative)
    aspect = image.width / max(image.height, 1)
    if aspect > 10 or aspect < 0.1:
        return False
    return True

def _ocr_image(image_bytes):
    """Perform OCR on a single image with optimized config."""
    try:
        image = Image.open(io.BytesIO(image_bytes))
        
        # Skip small/decorative images
        if not _should_ocr_image(image):
            return ""
        
        # Convert to RGB if needed (Tesseract works best with RGB)
        if image.mode not in ('L', 'RGB'):
            image = image.convert('RGB')
        
        # Use optimized config for faster block text detection
        text = pytesseract.image_to_string(image, config=OCR_CONFIG)
        return text.strip()
    except Exception:
        return ""

def _parallel_ocr(image_blobs, max_workers=4):
    """Run OCR on multiple images in parallel."""
    results = []
    
    if not image_blobs:
        return results
    
    with ThreadPoolExecutor(max_workers=max_workers) as executor:
        futures = {executor.submit(_ocr_image, blob): i for i, blob in enumerate(image_blobs)}
        for future in as_completed(futures):
            try:
                text = future.result()
                if text:
                    results.append(text)
            except:
                pass
    
    return results


# ============= FILE TYPE DETECTION =============

def detect_file_type(file_path):
    """
    Detects if file is PPTX or PDF based on extension or magic numbers.
    """
    # Simple extension check first
    if file_path.endswith(".pptx"):
        return "pptx"
    elif file_path.endswith(".pdf"):
        return "pdf"
    
    # Magic number check (MVP: just try opening)
    try:
        with open(file_path, 'rb') as f:
            header = f.read(4)
            if header == b'%PDF':
                return "pdf"
            # ZIP header for PPTX (PK..)
            elif header.startswith(b'PK'):
                return "pptx"
    except:
        pass
    return "unknown"


# ============= EXTRACTION FUNCTIONS =============

def extract_pptx(file_path):
    """
    Extracts text and runs OCR on images in a PPTX file, organized by slide.
    Optimized with parallel OCR and image filtering.
    """
    full_content = []
    
    try:
        prs = Presentation(file_path)
        
        for i, slide in enumerate(prs.slides):
            slide_text = []
            image_blobs = []
            
            # Text extraction
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
                
                # Collect images for parallel OCR
                if hasattr(shape, "image"):
                    try:
                        image_blobs.append(shape.image.blob)
                    except Exception:
                        pass
            
            # Parallel OCR on collected images
            slide_ocr = _parallel_ocr(image_blobs, max_workers=2)
            
            content_block = f"--- SLIDE {i+1} ---\n"
            if slide_text:
                content_block += "TEXT:\n" + "\n".join(slide_text) + "\n"
            if slide_ocr:
                content_block += "IMAGE OCR:\n" + "\n".join(slide_ocr) + "\n"
            
            full_content.append(content_block)
            
    except Exception as e:
        return f"Error reading PPTX: {str(e)}"

    return "\n".join(full_content)

def extract_pdf(file_path):
    """
    Extracts text and runs OCR on images in a PDF file, organized by page.
    Optimized with parallel OCR and image filtering.
    """
    full_content = []
    
    try:
        # Open with both pdfplumber and fitz
        pdf_plumber = pdfplumber.open(file_path)
        doc_fitz = fitz.open(file_path)
        
        for i in range(len(doc_fitz)):
            page_text = []
            image_blobs = []
            
            # Text extraction (pdfplumber)
            if i < len(pdf_plumber.pages):
                extracted = pdf_plumber.pages[i].extract_text()
                if extracted:
                    page_text.append(extracted.strip())
            
            # Image extraction (fitz) - collect for parallel OCR
            page = doc_fitz[i]
            for img in page.get_images(full=True):
                xref = img[0]
                try:
                    base_image = doc_fitz.extract_image(xref)
                    image_blobs.append(base_image["image"])
                except:
                    pass
            
            # Parallel OCR on collected images
            page_ocr = _parallel_ocr(image_blobs, max_workers=2)
            
            content_block = f"--- PAGE {i+1} ---\n"
            if page_text:
                content_block += "TEXT:\n" + "\n".join(page_text) + "\n"
            if page_ocr:
                content_block += "IMAGE OCR:\n" + "\n".join(page_ocr) + "\n"
            
            full_content.append(content_block)
            
        pdf_plumber.close()
        doc_fitz.close()
            
    except Exception as e:
        return f"Error reading PDF: {str(e)}"
        
    return "\n".join(full_content)

def process_file(file_path):
    """
    Main entry point to extract content from a file.
    Now includes caching for faster re-processing.
    """
    # Check cache first
    cached = _get_cached_extraction(file_path)
    if cached:
        return cached
    
    # Detect type if extension is missing (from downloader)
    ftype = detect_file_type(file_path)
    
    print(f"Processing {file_path} as {ftype}...")
    
    if ftype == "pptx":
        content = extract_pptx(file_path)
    elif ftype == "pdf":
        content = extract_pdf(file_path)
    else:
        # Try rename to .pptx and see? No, safe fail.
        # Check if it works as PPTX regardless of name
        try:
            content = extract_pptx(file_path)
        except:
            content = ""
    
    # Save to cache
    if content and not content.startswith("Error"):
        _save_extraction_cache(file_path, content)
    
    return content
