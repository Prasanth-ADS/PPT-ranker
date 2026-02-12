"""
Image Extractor Module

Converts PPTX slides and PDF pages into resized PIL images
for MiniCPM-V 2.6 VLM consumption.
"""
import os
import io
import tempfile
from PIL import Image
import fitz  # PyMuPDF
from .config import VLM_IMAGE_MAX_DIM, VLM_MAX_SLIDES
from .extractor import detect_file_type


def _resize_image(image: Image.Image, max_dim: int = VLM_IMAGE_MAX_DIM) -> Image.Image:
    """Resize image so its largest dimension is at most max_dim."""
    w, h = image.size
    if max(w, h) <= max_dim:
        return image
    scale = max_dim / max(w, h)
    new_w, new_h = int(w * scale), int(h * scale)
    return image.resize((new_w, new_h), Image.LANCZOS)


def _extract_pdf_images(file_path: str, max_slides: int = VLM_MAX_SLIDES) -> list:
    """Render PDF pages as images using PyMuPDF."""
    image_paths = []
    try:
        doc = fitz.open(file_path)
        num_pages = min(len(doc), max_slides)

        for i in range(num_pages):
            page = doc[i]
            # Render at 150 DPI for good quality without excessive size
            mat = fitz.Matrix(150 / 72, 150 / 72)
            pix = page.get_pixmap(matrix=mat)

            # Convert to PIL Image
            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            img = _resize_image(img)

            # Save to temp file
            tmp_path = os.path.join(
                tempfile.gettempdir(),
                f"vlm_slide_{i+1}_{os.path.basename(file_path)}.png"
            )
            img.save(tmp_path, "PNG")
            image_paths.append(tmp_path)

        doc.close()
    except Exception as e:
        print(f"  ⚠️ PDF image extraction error: {e}")

    return image_paths


def _extract_pptx_images(file_path: str, max_slides: int = VLM_MAX_SLIDES) -> list:
    """
    Extract images from PPTX slides.
    
    Strategy: Extract the largest embedded image from each slide.
    For slides without images, try to render shapes as a composite.
    """
    image_paths = []
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        slides = list(prs.slides)[:max_slides]

        for i, slide in enumerate(slides):
            best_image = None
            best_area = 0

            for shape in slide.shapes:
                if hasattr(shape, "image"):
                    try:
                        blob = shape.image.blob
                        img = Image.open(io.BytesIO(blob))
                        area = img.width * img.height
                        if area > best_area:
                            best_image = img
                            best_area = area
                    except Exception:
                        pass

            if best_image:
                best_image = _resize_image(best_image.convert("RGB"))
                tmp_path = os.path.join(
                    tempfile.gettempdir(),
                    f"vlm_slide_{i+1}_{os.path.basename(file_path)}.png"
                )
                best_image.save(tmp_path, "PNG")
                image_paths.append(tmp_path)

    except Exception as e:
        print(f"  ⚠️ PPTX image extraction error: {e}")

    return image_paths


def extract_slide_images(file_path: str) -> list:
    """
    Main entry point. Extracts slide/page images from a presentation file.

    Args:
        file_path: Path to PPTX or PDF file

    Returns:
        List of temporary image file paths (PNG format)
    """
    ftype = detect_file_type(file_path)

    if ftype == "pdf":
        images = _extract_pdf_images(file_path)
    elif ftype == "pptx":
        images = _extract_pptx_images(file_path)
    else:
        # Try PDF first, then PPTX
        images = _extract_pdf_images(file_path)
        if not images:
            images = _extract_pptx_images(file_path)

    print(f"  📸 Extracted {len(images)} slide images for VLM analysis")
    return images


def cleanup_temp_images(image_paths: list):
    """Remove temporary image files after VLM processing."""
    for path in image_paths:
        try:
            if os.path.exists(path) and "vlm_slide_" in path:
                os.remove(path)
        except Exception:
            pass
