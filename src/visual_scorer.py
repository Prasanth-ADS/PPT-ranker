import os
import re
import statistics
from pptx import Presentation
import pdfplumber
from collections import Counter
import numpy as np
# import cv2 # Optional if we do image processing, but let's stick to metadata for speed/reliability first

# Configuration
REQUIRED_KEYWORDS = {
    "problem": ["problem", "challenge", "background", "pain point", "context"],
    "solution": ["solution", "approach", "how it works", "implementation", "concept"],
    "architecture": ["architecture", "tech stack", "backend", "frontend", "diagram", "workflow"],
    "result": ["result", "outcome", "demo", "output", "testing", "validation"],
    "future": ["future", "roadmap", "next steps", "scalability"],
    "conclusion": ["conclusion", "summary", "closing"],
    "impact": ["impact", "value", "business", "market", "social"],
    "demo": ["demo", "prototype", "video", "screenshot"]
}
IDEAL_SLIDE_COUNT_RANGE = (8, 20)
IDEAL_WORD_COUNT_PER_SLIDE = (15, 120) # Slightly more flexible range

class VisualScorer:
    def __init__(self):
        pass

    def evaluate(self, file_path, file_type):
        """
        Main entry point. Returns a dictionary with:
        - visual_score (0-10 normalized, or user specified scale)
        - metrics (dict)
        - feedback (string)
        - format_compliant (bool)
        """
        if file_type == "pptx":
            return self._analyze_pptx(file_path)
        elif file_type == "pdf":
            return self._analyze_pdf(file_path)
        else:
            return self._empty_result("Unsupported format")

    def _empty_result(self, reason):
        return {
            "visual_score": 0,
            "metrics": {},
            "feedback": reason,
            "format_compliant": False
        }

    def _analyze_pptx(self, path):
        try:
            prs = Presentation(path)
            slides = prs.slides
            num_slides = len(slides)
            
            # 1. Section Compliance
            all_text = []
            slide_texts = []
            
            font_sizes = []
            font_names = []
            num_images = 0
            
            for slide in slides:
                s_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text = shape.text.strip()
                        if text:
                            s_text.append(text)
                            # Font extraction (heuristics)
                            if hasattr(shape, "text_frame"):
                                for paragraph in shape.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        if run.font.size:
                                            font_sizes.append(run.font.size.pt)
                                        if run.font.name:
                                            font_names.append(run.font.name)
                                            
                    # Image counting (Picture or GraphicFrame)
                    if shape.shape_type == 13: # PICTURE
                         num_images += 1
                    # GraphicFrames can be charts etc
                    if shape.shape_type == 7: # CHART
                         num_images += 0.5 # Partial credit for charts
                
                joined_s_text = " ".join(s_text).lower()
                slide_texts.append(joined_s_text)
                all_text.append(joined_s_text)

            # --- Scoring Logic ---
            
            # A. Section Compliance
            found_sections = []
            for category, aliases in REQUIRED_KEYWORDS.items():
                if any(any(alias in txt for alias in aliases) for txt in slide_texts):
                    found_sections.append(category)
            
            section_score = min(len(found_sections) / 5, 1.0) * 10 # 5 categories is excellent
            
            slide_count_score = 10 if (IDEAL_SLIDE_COUNT_RANGE[0] <= num_slides <= IDEAL_SLIDE_COUNT_RANGE[1]) else 5
            if num_slides < 3: slide_count_score = 0
            
            # B. Readability / Density (Max 30%)
            word_counts = [len(t.split()) for t in slide_texts]
            avg_words = statistics.mean(word_counts) if word_counts else 0
            if IDEAL_WORD_COUNT_PER_SLIDE[0] <= avg_words <= IDEAL_WORD_COUNT_PER_SLIDE[1]:
                density_score = 10
            else:
                density_score = 5 # Penalize extremes
                
            # C. Visual Elements (Max 20%)
            # Avg images per slide
            avg_imgs = num_images / num_slides if num_slides else 0
            visual_element_score = min(avg_imgs * 5, 10) # 2 images per slide = 10 points
            
            # D. Consistency (Max 20%)
            # Unique Fonts
            unique_fonts = len(set(font_names))
            if unique_fonts <= 3: consistency_score = 10
            elif unique_fonts <= 5: consistency_score = 7
            else: consistency_score = 4
            
            # -- Aggregation --
            # Weighted Sum to 100 then scale to 15
            # Sections: 30%, Density: 20%, Structure(Count): 10%, Visuals: 20%, Consistency: 20%
            
            # Fix scores to be 0-10 base
            final_raw = (
                (section_score * 3.0) +      # 0-30
                (density_score * 2.0) +      # 0-20
                (slide_count_score * 1.0) +  # 0-10
                (visual_element_score * 2.0) + # 0-20
                (consistency_score * 2.0)    # 0-20
            ) 
            # Total max = 100
            
            scaled_score = (final_raw / 100) * 15 # Scale to 15
            
            # Feedback Generation
            feedback_parts = []
            if section_score < 20: feedback_parts.append("Missing key sections (Problem/Solution/Arch).")
            if density_score < 10: feedback_parts.append(f"Text density is {'too high' if avg_words > 150 else 'too low'} ({int(avg_words)} words/slide).")
            if visual_element_score < 5: feedback_parts.append("Few visual elements/images detected.")
            if slide_count_score < 10: feedback_parts.append(f"Slide count ({num_slides}) is outside ideal range.")
            
            feedback = " ".join(feedback_parts) if feedback_parts else "Good format and visual structure."
            
            return {
                "visual_score": round(scaled_score, 1),
                "metrics": {
                    "slide_count": num_slides,
                    "avg_words": round(avg_words),
                    "images": num_images
                },
                "feedback": feedback,
                "format_compliant": num_slides >= 5 and len(found_sections) >= 2
            }
            
        except Exception as e:
            return self._empty_result(f"PPTX Error: {e}")

    def _analyze_pdf(self, path):
        try:
            with pdfplumber.open(path) as pdf:
                pages = pdf.pages
                num_pages = len(pages)
                
                all_text = []
                num_images = 0
                
                for page in pages:
                    text = page.extract_text() or ""
                    all_text.append(text.lower())
                    
                    # Count images
                    num_images += len(page.images)
                    
                # A. Sections
                found_sections = []
                for category, aliases in REQUIRED_KEYWORDS.items():
                    if any(any(alias in txt for alias in aliases) for txt in all_text):
                        found_sections.append(category)
                section_score = min(len(found_sections) / 5, 1.0) * 10
                
                # B. Density
                word_counts = [len(t.split()) for t in all_text]
                avg_words = statistics.mean(word_counts) if word_counts else 0
                density_score = 10 if (IDEAL_WORD_COUNT_PER_SLIDE[0] <= avg_words <= IDEAL_WORD_COUNT_PER_SLIDE[1]) else 5

                # C. Visuals
                avg_imgs = num_images / num_pages if num_pages else 0
                visual_element_score = min(avg_imgs * 5, 10)
                
                # D. Structure
                count_score = 10 if (IDEAL_SLIDE_COUNT_RANGE[0] <= num_pages <= IDEAL_SLIDE_COUNT_RANGE[1]) else 5
                
                # Consistency (Harder in PDF without raw font metadata easily) -> Default middle
                consistency_score = 8 
                
                # Calculate
                final_raw = (section_score * 3.0) + (density_score * 2.0) + (count_score * 1.0) + (visual_element_score * 2.0) + (consistency_score * 2.0)
                scaled_score = (final_raw / 100) * 15
                
                feedback_parts = []
                if section_score < 20: feedback_parts.append("Missing key sections.")
                if density_score < 10: feedback_parts.append(f"Text density issues ({int(avg_words)} words/page).")
                
                feedback = " ".join(feedback_parts) if feedback_parts else "Good PDF structure."

                return {
                    "visual_score": round(scaled_score, 1),
                    "metrics": {
                        "slide_count": num_pages,
                        "avg_words": round(avg_words),
                        "images": num_images
                    },
                    "feedback": feedback,
                    "format_compliant": num_pages >= 5
                }
        except Exception as e:
            return self._empty_result(f"PDF Error: {e}")

